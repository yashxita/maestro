import pdfplumber
from pptx import Presentation
import fitz
from PIL import Image
import pytesseract
import io
import os

def extract_text_from_pdf(path: str) -> str:
    text_parts = []

    # Try direct text extraction first
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            txt = page.extract_text()
            if txt:
                text_parts.append(txt)

    # OCR fallback
    pdf_doc = fitz.open(path)
    for page_index, page in enumerate(pdf_doc):
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = pdf_doc.extract_image(xref)
            image_bytes = base_image.get("image")
            image_ext = base_image.get("ext", "")
            if not image_bytes:
                continue

            try:
                # safer open via BytesIO
                image = Image.open(io.BytesIO(image_bytes))
                ocr_text = pytesseract.image_to_string(image)
                if ocr_text.strip():
                    text_parts.append(ocr_text)
            except Exception as e:
                print(f"Skipping image on page {page_index+1}, index {img_index}: {e}")
                continue

    return "\n".join(text_parts)


def extract_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    text_parts = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text)

    # OCR for images
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # picture
                img_bytes = shape.image.blob
                image = Image.open(io.BytesIO(img_bytes))
                ocr_text = pytesseract.image_to_string(image)
                if ocr_text.strip():
                    text_parts.append(ocr_text)

    return "\n".join(text_parts)


def extract_text(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(path)
    elif ext == ".pptx":
        return extract_text_from_pptx(path)
    else:
        raise ValueError("Unsupported file type")
